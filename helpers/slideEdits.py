import collections
import collections.abc
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm
from pptx.util import Pt


# This is for the first line in the program
# it simply copies the current line to the editedText variable
#
def firstPartitionEdit(editedText, curLine):
    editedText = curLine
    return editedText


# This is for the lines appended to the end of a verse line
# For instance if we have the verse
#     Sunt numai un vas de lut
#     dar de Dumnezeu facut
#     sunt creat si modelat
#     doar de Dumnezeu preananlt
# Then this function will deal with the SECOND and FOURTH lines with the result being
#     Sunt numai un vas de lut dar de Dumnezeu facut
#     sunt creat si modelat doar de Dumnezeu preanalt
def middlePartitionEdit(editedText, curLine):
    if (editedText[-1] in [".", ",", "!", "?", ":", ";"]):
        if (curLine[-1] in [".", ",", "!", "?", ":", ";"]):
            editedText = editedText[:-1] + ", " + curLine[:-1]
        else:
            editedText = editedText[:-1] + ", " + curLine

    else:
        if (curLine[-1] in [".", ",", "!", "?", ":", ";"]):
            editedText = editedText + ", " + curLine[:-1]
        else:
            editedText = editedText + ", " + curLine

    return editedText


# This is new lines
# For instance if we have the verse
#     Sunt numai un vas de lut
#     dar de Dumnezeu facut
#     sunt creat si modelat
#     doar de Dumnezeu preananlt
# Then this function will deal with the THIRD line with the result being
#     Sunt numai un vas de lut dar de Dumnezeu facut
#     sunt creat si modelat doar de Dumnezeu preanalt
def newLinePartitionEdit(editedText, curLine):
    if (editedText[-1] in [".", ",", "!", "?", ":", ";"]):
        editedText = editedText[:-1] + '\n' + curLine
    else:
        editedText = editedText + '\n' + curLine

    return editedText


# Takes the text and puts it on 2 lines
# d
#
def editText(filename, shape):
    lines = shape.text.splitlines(False)

    for line in lines:
        if (line == '' or line.isspace() == True):
            lines.remove(line)

    # Leave these here for future testing purposes
    # print(lines)
    # print(len(lines))

    # This slide is either empty or already formatted properlys
    # if (len(lines) < 2):
    #     return False

    else:
        # This makes it so that all text is made to be 2 lines. See the for loop below
        addNewLine = len(lines) / 2
        editedText = ''
        isFirstLine = True

        for i in range(len(lines)):
            curLine = lines[i].rstrip()

            if (len(curLine) == 0):
                print(
                    "WARNING - there is a stray empty line in the slide, ignoring it...")
                continue

            if (isFirstLine):
                # print("i == 0 lines: " +
                #       str(lines[i]) + "edited text: " + str(editedText) + "\n\n")

                editedText = firstPartitionEdit(editedText, curLine)
                isFirstLine = False
            elif (i % addNewLine == 0):
                # print("i mod line == 0 lines: " +
                #       str(curLine) + "edited text: " + str(editedText) + "\n\n")

                editedText = newLinePartitionEdit(editedText, curLine)
            else:
                # print("else == 0 lines: " +
                #       str(curLine) + "edited text: " + str(editedText) + "\n\n")

                curLine = curLine[0].lower() + curLine[1:]
                editedText = middlePartitionEdit(editedText, curLine)

        text_frame = shape.text_frame
        text_frame.clear()  # removes all paragraphs except for one empty one
        p = text_frame.paragraphs[0]  # access the empty paragraph
        run = p.add_run()  # adds a run to the new empty paragraph
        run.text = editedText  # change the text of the paragraph to something new
        return True


# Applies formatting to the characters of the paragraphs in a text frame
# Changes things like font size, font family, etc...
# NOTE: font color should be changed in the paragraph formatting level
def editCharacterFormatting(filename, shape, textValuesDict):
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.runs[0]
    font = run.font

    try:
        font.name = textValuesDict['fontName']
    except Exception as e:
        print("WARNING | " + filename + ": " + str(e) +
              "--> something went wrong, using default font family value of Calibri")
        font.name = 'Calibri'

    try:
        font.size = Pt(int(textValuesDict['fontSize']))
    except Exception as e:
        print("WARNING | " + filename + ": " + str(e) +
              "--> something went wrong, using default font size value of 20")
        font.size = Pt(20)

    try:
        font.bold = textValuesDict['fontBold']
    except Exception as e:
        print("WARNING | " + filename + ": " + str(e) +
              "--> something went wrong, using default font bold value of True")
        font.bold = True

    try:
        font.italic = textValuesDict['fontItalic']
    except Exception as e:
        print("WARNING | " + filename + ": " + str(e) +
              "--> something went wrong, using default font family value of False")
        font.italic = False


# Applies formatting to the paragraphs of the shapes with text_frames in the slide
# Should not change character formatting other than text color
# NOTE: There should only be one paragraph in the text frame at this point
def editParagraphFormatting(filename, shape, textValuesDict):
    text_frame = shape.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


# Edits other shapes on the slide
# Edits text using desired values
# Gets rid of pictures
def editShape(filename, shape, textValuesDict):
    if shape.has_text_frame:
        shape.width = Cm(33.867)
        shape.height = Cm(19.05)
        shape.left = Cm(0)
        shape.top = Cm(0)

        success = editText(filename, shape)
        if success == True:
            editParagraphFormatting(filename, shape, textValuesDict)
            editCharacterFormatting(filename, shape, textValuesDict)
        else:
            shape._element.getparent().remove(shape._element)

    else:
        shape._element.getparent().remove(shape._element)
