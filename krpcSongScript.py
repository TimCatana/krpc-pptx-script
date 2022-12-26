import os
import sys
import collections 
import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm

from helpers.inputs import getDesiredInputs
from helpers.readCSVFile import readCSVFile
from helpers.slideEdits import editShape

def main():

  # Notify user to close all powerpoint slides
  input("IMPORTANT Before beginning please ensure all instances of powerpoint are closed and that your input.csv file is in the correct form (example: 21,Calibri,y,n ). Press Enter to begin... ") 

  # Make the edited directory if it doesn't exist
  try:
    os.mkdir('edited')
    print("NOTICE: edited directory created, continuing...")
  except OSError:
    print("NOTICE: edited directory already exists, continuing...")

  # Read the csv file and get the contents
  textValuesDict = readCSVFile('input.csv')

  # Get the powerpoint slides in the root directory and so the edit for each one
  for filename in os.listdir("./"):

    if (filename[-4:] == ".ppt" or filename[-5:] == ".pptx"):

      try:
        prs = Presentation(filename)
      except BaseException as e:
        print("ERROR: " + str(e))
        print("WARNING: Failed to open powerpoint file: " + filename + " continuing...")
        continue
      
      prs.slide_width = Cm(33.867)
      prs.slide_height = Cm(19.05)

      for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

        for shape in slide.shapes:
          editShape(filename, shape, textValuesDict)

      try:
        prs.save('edited/' + filename.rsplit('.', 1)[0] + '.pptx')
        print("Finished Successfully")
      except BaseException as error:
        print('ERROR STOPPING SCRIPT: Failed to save file: ' + '\'edited/' + filename.split('.')[0] + '.pptx\' ' + 'please make sure that all instances of powerpoint are closed')
        input("PRESS ENTER TO CLOSE TERMINAL") 
        sys.exit()

      input("PRESS ENTER TO CLOSE TERMINAL") 
      

if __name__ == "__main__":
  main()