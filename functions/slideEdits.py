import collections 
import collections.abc
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Pt

# Takes the text and puts it on 2 lines 
#
# 
def editText(filename, shape):
  lines = shape.text.splitlines(False)
  
  for line in lines:
    if(line == '' or line.isspace() == True):
      lines.remove(line)
    # TODO - get rid of whitespace on the end of the strings, strip() was not working when I tried implementing it
  
  # Leave these here for future testing purposes
  # print(lines)
  # print(len(lines))
  
  if(len(lines) < 2):
    return False
 
  else:
    addNewLine = len(lines) / 2 # This makes it so that all text is made to be 2 lines. See the for loop below
    editedText = ''

    for i in range(len(lines)):
      lines[i].rstrip()
      if(i == 0):
        editedText = lines[i] + ' '
      elif(i % addNewLine == 0):
        editedText = editedText + '\n' + lines[i] + ' '
      else:
        editedText = editedText + lines[i] + ' '

    text_frame = shape.text_frame
    text_frame.clear() # removes all paragraphs except for one empty one
    p = text_frame.paragraphs[0] # access the empty paragraph
    run = p.add_run() # adds a run to the new empty paragraph
    run.text = editedText # change the text of the paragraph to something new
    return True

# Applies formatting to the characters of the paragraphs in a text frame
# Changes things like font size, font family, etc...
# NOTE: font color should be changed in the paragraph formatting level
def editCharacterFormatting(filename, shape, textValuesDict):
  text_frame = shape.text_frame
  p = text_frame.paragraphs[0]
  run = p.runs[0]
  font = run.font

  try:
    font.name = textValuesDict['fontName']
  except Exception as e:
    print("WARNING | " + filename + ": " + str(e) + "--> something went wrong, using default font family value of Calibri")
    font.name = 'Calibri'

  try:
    font.size = Pt(int(textValuesDict['fontSize']))
  except Exception as e:
    print("WARNING | " + filename + ": " + str(e) + "--> something went wrong, using default font size value of 20") 
    font.size = Pt(20)
  
  try:
    font.bold = textValuesDict['fontBold']
  except Exception as e:
    print("WARNING | " + filename + ": " + str(e) + "--> something went wrong, using default font bold value of True")
    font.bold = True

  try:
    font.italic = textValuesDict['fontItalic']
  except Exception as e:
    print("WARNING | " + filename + ": " + str(e) + "--> something went wrong, using default font family value of False")
    font.italic = False

# Applies formatting to the paragraphs of the shapes with text_frames in the slide
# Should not change character formatting other than text color
def editParagraphFormatting(filename, shape, textValuesDict):
  text_frame = shape.text_frame
  
  for p in text_frame.paragraphs:
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

# Edits other shapes on the slide
# Edits text using desired values
# Gets rid of pictures
def editShape(filename, shape, textValuesDict):
  if shape.has_text_frame:
    shape.width = Cm(33.867)
    shape.height = Cm(19.05)
    shape.left = Cm(0)
    shape.top = Cm(0)

    success = editText(filename, shape)
    if success == True:
      editParagraphFormatting(filename, shape, textValuesDict)
      editCharacterFormatting(filename, shape, textValuesDict)
    else:
      shape._element.getparent().remove(shape._element)
    
  else :  
    shape._element.getparent().remove(shape._element)